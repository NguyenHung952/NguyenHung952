from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from services.design_engine import slide_theme

ROLE_FONT = {
    "title": 42,
    "introduction": 32,
    "section": 34,
    "key_point": 30,
    "content": 28,
    "summary": 30,
}


def _layout_for(prs: Presentation, layout: str):
    if layout == "title":
        return prs.slide_layouts[0]
    if layout == "section":
        return prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[1]
    return prs.slide_layouts[1]


def build_pptx(slides: list[dict]) -> bytes:
    prs = Presentation()
    accent = RGBColor(52, 89, 149)

    for idx, slide_data in enumerate(slides):
        requested_layout = slide_data.get("layout", "content")
        role = slide_data.get("role", "content")
        theme = slide_theme(role)
        margin_left = Inches(theme["margin_left"])

        layout = _layout_for(prs, requested_layout if idx > 0 else "title")
        slide = prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        if title_shape:
            title_shape.left = margin_left
            title_shape.text = slide_data.get("title", "Untitled")
            title_p = title_shape.text_frame.paragraphs[0]
            title_p.font.size = Pt(theme["title_size"] if role != "title" else ROLE_FONT["title"])
            title_p.font.bold = True
            title_p.font.color.rgb = accent if role == "key_point" else RGBColor(30, 30, 30)
            title_p.alignment = PP_ALIGN.LEFT

        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            body.left = margin_left
            tf = body.text_frame
            tf.clear()
            subtitle = slide_data.get("subtitle") or ""
            if subtitle:
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(18)
                p.font.italic = True
                p.space_after = Pt(8)
            for bullet in (slide_data.get("bullets") or [])[:5]:
                p = tf.add_paragraph()
                p.text = f"• {bullet}"
                p.level = 0
                p.space_after = Pt(theme["bullet_spacing"])
                p.font.size = Pt(20 if role == "key_point" else 18)

        image_url = slide_data.get("image_url") or slide_data.get("image")
        if image_url:
            width = Inches(4.6)
            height = Inches(4.8)
            left = Inches(0.6) if requested_layout == "image_left" else Inches(8.2)
            top = Inches(1.5)
            box = slide.shapes.add_textbox(left, top, width, height)
            box.text_frame.text = f"Image: {image_url}"
            box.text_frame.word_wrap = True
            box.text_frame.paragraphs[0].font.size = Pt(11)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
